"""Build the Watcha submission .pptx for Resona.

Event rules mandate Alibaba PuHui 3.0 + Times New Roman. Our live stage
deck stays in deck/index.html with its own typography; this artifact
is for the submission packet only.

Usage:
    python scripts/build_submission_pptx.py
    # writes submission/Resona.pptx

Layout system:
    - 16:9 widescreen, 13.333 x 7.5 in (EMU math handled by python-pptx)
    - Dark ink background, bone text, brass accents
    - PuHui 3.0 for headlines + body
    - Times New Roman for big numbers, ratios, citations, and "serif moments"
"""

from __future__ import annotations

import io
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- Brand tokens (matched to deck/index.html :root vars) ---------------------

INK_0 = RGBColor(0x0A, 0x0B, 0x10)
INK_1 = RGBColor(0x12, 0x13, 0x1A)
INK_2 = RGBColor(0x1A, 0x1C, 0x26)
INK_3 = RGBColor(0x24, 0x27, 0x32)

BONE_0 = RGBColor(0xF4, 0xEC, 0xE1)
BONE_1 = RGBColor(0xE4, 0xD9, 0xC4)
BONE_2 = RGBColor(0xB8, 0xAC, 0x94)
BONE_3 = RGBColor(0x8A, 0x82, 0x72)

BRASS = RGBColor(0xC9, 0xA9, 0x6E)
BRASS_BRIGHT = RGBColor(0xE6, 0xC6, 0x8A)
BRASS_DIM = RGBColor(0x7A, 0x65, 0x42)

PULSE = RGBColor(0x7B, 0xC1, 0x96)
WARN = RGBColor(0xD1, 0x85, 0x89)

FONT_SANS = "Alibaba PuHui 3.0"
FONT_SERIF = "Times New Roman"

# --- Geometry -----------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)


# --- Low-level helpers --------------------------------------------------------


def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None,
             line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    return shape


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = 0
    tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    return tb


def write_runs(paragraph, runs, default_font=FONT_SANS):
    """Write a sequence of (text, opts) tuples as runs on an existing paragraph.

    opts keys: font (name), size (Pt), color (RGBColor), bold, italic,
    letter_spacing (0.01 = 1%), tracking in Pt.
    """
    first = True
    for text, opts in runs:
        r = paragraph.add_run() if not first else paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        first = False
        r.text = text
        r.font.name = opts.get("font", default_font)
        if "size" in opts:
            r.font.size = Pt(opts["size"])
        if "color" in opts:
            r.font.color.rgb = opts["color"]
        if opts.get("bold"):
            r.font.bold = True
        if opts.get("italic"):
            r.font.italic = True


def set_text(paragraph, text, *, font=FONT_SANS, size=14, color=BONE_0,
             bold=False, italic=False, align=None):
    if align is not None:
        paragraph.alignment = align
    if paragraph.runs:
        r = paragraph.runs[0]
    else:
        r = paragraph.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic


def para(textframe, text="", **kwargs):
    """Append a fresh paragraph to a textframe and style it."""
    p = textframe.add_paragraph()
    if text:
        set_text(p, text, **kwargs)
    elif kwargs.get("align") is not None:
        p.alignment = kwargs["align"]
    return p


def first_para(textframe, text="", **kwargs):
    """Replace the implicit first paragraph's text (textframe always starts
    with one empty paragraph)."""
    p = textframe.paragraphs[0]
    set_text(p, text, **kwargs)
    return p


# --- Slide chrome helpers -----------------------------------------------------


def add_chrome_top(slide, case_label: str, index: str):
    """Top strip with 'Case 0N · Title' left and brand + index right."""
    y = Inches(0.35)
    # Left: case label
    tb = add_textbox(slide, MARGIN, y, Inches(6), Inches(0.35))
    set_text(tb.text_frame.paragraphs[0], case_label,
             size=11, color=BRASS, bold=True, font=FONT_SANS)
    # Right: brand + index
    tb = add_textbox(slide, SLIDE_W - MARGIN - Inches(4), y, Inches(4), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r1 = p.add_run()
    r1.text = "Resona   "
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(16)
    r1.font.color.rgb = BONE_0
    r2 = p.add_run()
    r2.text = index
    r2.font.name = FONT_SANS
    r2.font.size = Pt(10)
    r2.font.color.rgb = BONE_3
    r2.font.bold = True
    # Hairline under chrome
    add_rect(slide, MARGIN, Inches(0.75), SLIDE_W - 2 * MARGIN, Emu(9525),
             fill_color=BONE_3)


def add_chrome_bottom(slide, footer_text: str, page: str):
    y = SLIDE_H - Inches(0.55)
    # Hairline
    add_rect(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Emu(9525), fill_color=BONE_3)
    # Left text
    tb = add_textbox(slide, MARGIN, y + Inches(0.1), Inches(8), Inches(0.35))
    set_text(tb.text_frame.paragraphs[0], footer_text,
             size=9, color=BONE_3, font=FONT_SANS, bold=True)
    # Right page
    tb = add_textbox(slide, SLIDE_W - MARGIN - Inches(2), y + Inches(0.1),
                     Inches(2), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    set_text(p, page, size=9, color=BONE_3, font=FONT_SANS, bold=True)


# --- Slide 01 · Hook ----------------------------------------------------------


def build_slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, INK_1)
    add_chrome_top(slide, "CASE 01 · HOOK", "01 / 05")

    # Left column
    left_x = MARGIN
    left_w = Inches(6.2)

    # Eyebrow
    tb = add_textbox(slide, left_x, Inches(1.3), left_w, Inches(0.4))
    set_text(tb.text_frame.paragraphs[0], "THE NUMBER YOU CANNOT SIT WITH",
             size=11, color=BRASS, bold=True, font=FONT_SANS)

    # Pre text
    tb = add_textbox(slide, left_x, Inches(1.75), left_w, Inches(0.6))
    set_text(tb.text_frame.paragraphs[0],
             "The average UK knowledge worker now spends",
             size=18, color=BONE_1, font=FONT_SANS)

    # Hero number
    tb = add_textbox(slide, left_x, Inches(2.3), left_w, Inches(2.2))
    p = tb.text_frame.paragraphs[0]
    r1 = p.add_run()
    r1.text = "9.5"
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(160)
    r1.font.color.rgb = BONE_0
    r1.font.bold = False

    # Unit line
    tb = add_textbox(slide, left_x, Inches(4.45), left_w, Inches(0.45))
    set_text(tb.text_frame.paragraphs[0], "HOURS PER DAY",
             size=12, color=BRASS, bold=True, font=FONT_SANS)

    # Post line
    tb = add_textbox(slide, left_x, Inches(4.85), left_w, Inches(0.45))
    set_text(tb.text_frame.paragraphs[0], "sitting at a desk.",
             size=20, color=BONE_1, font=FONT_SANS)

    # Callout
    tb = add_textbox(slide, left_x, Inches(5.4), left_w, Inches(0.8))
    set_text(tb.text_frame.paragraphs[0], "More than they sleep.",
             size=28, color=BRASS_BRIGHT, italic=True, font=FONT_SERIF)
    # Brass vertical bar next to callout
    add_rect(slide, left_x, Inches(5.4), Emu(25400), Inches(0.8),
             fill_color=BRASS)

    # Right column — research card
    right_x = Inches(7.4)
    right_w = Emu(SLIDE_W - right_x - MARGIN)

    # Day bar header
    tb = add_textbox(slide, right_x, Inches(1.3), right_w, Inches(0.4))
    set_text(tb.text_frame.paragraphs[0], "A WEEKDAY, SAMPLED",
             size=10, color=BONE_3, bold=True, font=FONT_SANS)

    # Day bar: four segments
    bar_top = Inches(1.75)
    bar_h = Inches(0.7)
    segments = [
        ("SLEEP 8h", 0.267, INK_3, BONE_2),          # 8 / 30 active-ish; display simplified
        ("SITTING 9.5h", 0.40, BRASS, INK_0),
        ("MOVE 2.5h", 0.105, RGBColor(0x3E, 0x50, 0x47), BONE_1),
        ("OTHER 4h", 0.228, INK_3, BONE_2),
    ]
    # Normalise widths to fill the card
    total_frac = sum(s[1] for s in segments)
    cursor = right_x
    for label, frac, fill, text in segments:
        w = Emu(int(right_w * (frac / total_frac)))
        add_rect(slide, cursor, bar_top, w, bar_h,
                 fill_color=fill, line_color=fill)
        tb = add_textbox(slide, cursor + Inches(0.1), bar_top + Inches(0.2),
                         w - Inches(0.2), bar_h - Inches(0.2))
        set_text(tb.text_frame.paragraphs[0], label,
                 size=9, color=text, bold=True, font=FONT_SANS)
        cursor += w

    # Source line
    tb = add_textbox(slide, right_x, Inches(2.55), right_w, Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r1 = p.add_run()
    r1.text = "ONS UK Time Use Survey · 2022-23"
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(10)
    r1.font.color.rgb = BONE_3
    r1.font.italic = True
    r2 = p.add_run()
    r2.text = "           40% SITTING"
    r2.font.name = FONT_SANS
    r2.font.size = Pt(10)
    r2.font.color.rgb = BRASS
    r2.font.bold = True

    # Research card
    card_top = Inches(3.15)
    card_h = Inches(3.25)
    add_rect(slide, right_x, card_top, right_w, card_h,
             fill_color=INK_2, line_color=BRASS, line_width=Pt(1))
    # Top brass hairline inside card
    add_rect(slide, right_x + Inches(0.3), card_top + Inches(0.04),
             right_w - Inches(0.6), Emu(12700), fill_color=BRASS_BRIGHT)

    facts = [
        ("CARDIOVASCULAR", "+147", "% RISK", "Wilmot et al, Diabetologia, 2012"),
        ("TYPE 2 DIABETES", "+112", "% RISK", "meta-analysis, 794,577 people"),
        ("ALL-CAUSE MORTALITY", "+40", "%", "Ekelund et al, The Lancet, 2016", "AT 8+ HRS/DAY"),
    ]
    col_w = Inches((right_w.inches - 0.6) / 3)
    fx = right_x + Inches(0.3)
    fy = card_top + Inches(0.35)

    for i, fact in enumerate(facts):
        label, num, unit, src = fact[0], fact[1], fact[2], fact[3]
        cond = fact[4] if len(fact) > 4 else None
        cx = fx + Inches(i * (col_w.inches))
        # Label
        tb = add_textbox(slide, cx, fy, col_w, Inches(0.35))
        set_text(tb.text_frame.paragraphs[0], label,
                 size=11, color=BONE_1, bold=True, font=FONT_SANS)
        # Big number (serif) + unit (sans)
        tb = add_textbox(slide, cx, fy + Inches(0.4), col_w, Inches(1.1))
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run()
        r1.text = num
        r1.font.name = FONT_SERIF
        r1.font.size = Pt(54)
        r1.font.color.rgb = BRASS_BRIGHT
        r2 = p.add_run()
        r2.text = unit
        r2.font.name = FONT_SANS
        r2.font.size = Pt(18)
        r2.font.color.rgb = BRASS
        r2.font.bold = True
        # Optional condition line
        if cond:
            tb = add_textbox(slide, cx, fy + Inches(1.5), col_w, Inches(0.3))
            set_text(tb.text_frame.paragraphs[0], cond,
                     size=9, color=BONE_2, bold=True, font=FONT_SANS)
        # Source
        tb = add_textbox(slide, cx, fy + Inches(1.85), col_w, Inches(0.45))
        set_text(tb.text_frame.paragraphs[0], src,
                 size=10, color=BONE_3, italic=True, font=FONT_SERIF)

    add_chrome_bottom(slide,
                      "SEDENTARY · CARDIOVASCULAR · METABOLIC · RESPIRATORY",
                      "01 / 05")


# --- Slide 02 · Problem -------------------------------------------------------


def build_slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, INK_2)
    add_chrome_top(slide, "CASE 02 · PROBLEM", "02 / 05")

    # Lead
    tb = add_textbox(slide, MARGIN, Inches(1.2),
                     SLIDE_W - 2 * MARGIN, Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p, "THE MARKET IS MOVING. THE TOOLS ARE NOT.",
             size=12, color=BRASS, bold=True, font=FONT_SANS, align=PP_ALIGN.CENTER)

    # Headline
    tb = add_textbox(slide, MARGIN, Inches(1.7),
                     SLIDE_W - 2 * MARGIN, Inches(1.4))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "Wall Street is "
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(48)
    r1.font.color.rgb = BONE_0
    r2 = p1.add_run()
    r2.text = "buying"
    r2.font.name = FONT_SERIF
    r2.font.size = Pt(48)
    r2.font.color.rgb = BRASS
    r2.font.italic = True
    r3 = p1.add_run()
    r3.text = " wellness."
    r3.font.name = FONT_SERIF
    r3.font.size = Pt(48)
    r3.font.color.rgb = BONE_0
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r4 = p2.add_run()
    r4.text = "The hardware is "
    r4.font.name = FONT_SERIF
    r4.font.size = Pt(48)
    r4.font.color.rgb = BONE_0
    r5 = p2.add_run()
    r5.text = "failing"
    r5.font.name = FONT_SERIF
    r5.font.size = Pt(48)
    r5.font.color.rgb = BRASS
    r5.font.italic = True
    r6 = p2.add_run()
    r6.text = " it."
    r6.font.name = FONT_SERIF
    r6.font.size = Pt(48)
    r6.font.color.rgb = BONE_0

    # Two case cards
    case_top = Inches(3.8)
    case_h = Inches(2.85)
    case_w = Inches((SLIDE_W.inches - 2 * MARGIN.inches - 0.5) / 2)
    gap = Inches(0.5)

    def draw_case(x, label, tag, tag_color, source, title_lines, body, footer):
        add_rect(slide, x, case_top, case_w, case_h,
                 fill_color=INK_1, line_color=BONE_3, line_width=Pt(0.5))
        # Top brass accent
        add_rect(slide, x, case_top, case_w, Pt(2), fill_color=BRASS)
        # Header row
        tb = add_textbox(slide, x + Inches(0.3), case_top + Inches(0.2),
                         case_w - Inches(0.6), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run()
        r1.text = label
        r1.font.name = FONT_SANS
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = BRASS
        r2 = p.add_run()
        r2.text = f"       {tag}"
        r2.font.name = FONT_SANS
        r2.font.size = Pt(9)
        r2.font.bold = True
        r2.font.color.rgb = tag_color
        # Source
        tb = add_textbox(slide, x + Inches(0.3), case_top + Inches(0.55),
                         case_w - Inches(0.6), Inches(0.3))
        set_text(tb.text_frame.paragraphs[0], source,
                 size=9, color=BONE_3, italic=True, font=FONT_SERIF)
        # Title
        tb = add_textbox(slide, x + Inches(0.3), case_top + Inches(0.9),
                         case_w - Inches(0.6), Inches(0.9))
        tf = tb.text_frame
        for i, line in enumerate(title_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            set_text(p, line, size=22, color=BONE_0, font=FONT_SERIF)
        # Body
        tb = add_textbox(slide, x + Inches(0.3), case_top + Inches(1.85),
                         case_w - Inches(0.6), Inches(0.75))
        set_text(tb.text_frame.paragraphs[0], body,
                 size=11, color=BONE_1, font=FONT_SANS)
        # Footer
        tb = add_textbox(slide, x + Inches(0.3), case_top + case_h - Inches(0.4),
                         case_w - Inches(0.6), Inches(0.3))
        set_text(tb.text_frame.paragraphs[0], footer,
                 size=9, color=BONE_3, italic=True, font=FONT_SERIF)

    draw_case(
        x=MARGIN,
        label="CASE A",
        tag="DEMAND VALIDATED",
        tag_color=PULSE,
        source="Hedgeweek · July 2025",
        title_lines=["Citadel hires its first", "Chief Medical Officer."],
        body="Ken Griffin appoints Dr David Stark (ex-Morgan Stanley, "
             "Harvard-trained pediatric neurologist) to drive peak "
             "performance across the firm. A signal the $5.7 trillion "
             "hedge fund industry treats employee health as alpha.",
        footer="First-ever CMO on Wall Street · July 2025",
    )
    draw_case(
        x=MARGIN + case_w + gap,
        label="CASE B",
        tag="SUPPLY BROKEN",
        tag_color=WARN,
        source="Healthcare Innovation · Fitbit corporate study",
        title_lines=["The wearable gap.", "£1,000 saved, then abandoned."],
        body="Enrolled employees saved an average £1,000 per year in "
             "healthcare costs but 40% dropped off by month three. Fitbit "
             "corporate sales have declined every year since 2018. "
             "The demand is there. The hardware is not the answer.",
        footer="Global corporate wellness market: $53 bn · mostly unloved",
    )

    # Foot
    tb = add_textbox(slide, MARGIN, SLIDE_H - Inches(1.1),
                     SLIDE_W - 2 * MARGIN, Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p, "No daily, private, hardware-free check-in for a whole team. Until now.",
             size=18, color=BONE_1, italic=True, font=FONT_SERIF,
             align=PP_ALIGN.CENTER)

    add_chrome_bottom(slide,
                      "Citadel · July 2025   ·   Fitbit corporate · Healthcare Innovation study",
                      "02 / 05")


# --- Slide 03 · Solution ------------------------------------------------------


def build_slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, INK_1)
    add_chrome_top(slide, "CASE 03 · SOLUTION", "03 / 05")

    # Eyebrow
    tb = add_textbox(slide, MARGIN, Inches(1.2),
                     SLIDE_W - 2 * MARGIN, Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p, "THREE SIGNALS · ONE PHONE · TWO MINUTES",
             size=12, color=BRASS, bold=True, font=FONT_SANS,
             align=PP_ALIGN.CENTER)

    # Headline
    tb = add_textbox(slide, MARGIN, Inches(1.7),
                     SLIDE_W - 2 * MARGIN, Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "Every phone is an "
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(54)
    r1.font.color.rgb = BONE_0
    r2 = p.add_run()
    r2.text = "instrument."
    r2.font.name = FONT_SERIF
    r2.font.size = Pt(54)
    r2.font.color.rgb = BRASS
    r2.font.italic = True

    # Three module cards
    card_top = Inches(3.5)
    card_h = Inches(3.0)
    card_w = Inches((SLIDE_W.inches - 2 * MARGIN.inches - 1.0) / 3)
    gap = Inches(0.5)

    modules = [
        ("01", "Breath.", "Microphone · Acoustic spirometry",
         "FEV1, FVC, PEF. Forced exhalation analysed from turbulent flow "
         "noise. Normalised to Hankinson NHANES III predicted values.",
         "6s blow · 48 kHz PCM · on-device", "LIVE", True),
        ("02", "Motion.", "IMU · DeviceMotion API",
         "Tremor from 10s stillness. Gait cadence and stride variability "
         "from 10-step walk. Standing-up streaks over time.",
         "60 Hz accelerometer · FFT bands", "LIVE", True),
        ("03", "Heart.", "Front camera · rPPG",
         "Resting heart rate and HRV from 30s face video. Sub-pixel colour "
         "oscillations, bandpass 0.7 to 4 Hz.",
         "MediaPipe ROI · Q3 2026", "Q3", False),
    ]

    for i, (n, title, sensor, signals, ref, status, is_live) in enumerate(modules):
        x = MARGIN + Inches(i * (card_w.inches + gap.inches))
        fill = INK_2 if is_live else INK_1
        border = BRASS if is_live else BONE_3
        add_rect(slide, x, card_top, card_w, card_h,
                 fill_color=fill, line_color=border,
                 line_width=Pt(1 if is_live else 0.5))
        # Header: number + status pill
        tb = add_textbox(slide, x + Inches(0.3), card_top + Inches(0.25),
                         card_w - Inches(0.6), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = n
        r.font.name = FONT_SERIF
        r.font.size = Pt(40)
        r.font.color.rgb = BRASS if is_live else BONE_3
        r.font.italic = True
        # Status pill on right
        status_color = PULSE if is_live else BRASS
        tb = add_textbox(slide,
                         x + card_w - Inches(1.0),
                         card_top + Inches(0.35),
                         Inches(0.8), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        set_text(p, status, size=10, color=status_color, bold=True,
                 font=FONT_SANS, align=PP_ALIGN.RIGHT)
        # Title
        tb = add_textbox(slide, x + Inches(0.3), card_top + Inches(1.05),
                         card_w - Inches(0.6), Inches(0.5))
        set_text(tb.text_frame.paragraphs[0], title,
                 size=28, color=BONE_0, font=FONT_SERIF)
        # Sensor
        tb = add_textbox(slide, x + Inches(0.3), card_top + Inches(1.55),
                         card_w - Inches(0.6), Inches(0.3))
        set_text(tb.text_frame.paragraphs[0], sensor,
                 size=10, color=BRASS, bold=True, font=FONT_SANS)
        # Signals
        tb = add_textbox(slide, x + Inches(0.3), card_top + Inches(1.9),
                         card_w - Inches(0.6), Inches(0.8))
        set_text(tb.text_frame.paragraphs[0], signals,
                 size=11, color=BONE_2, font=FONT_SANS)
        # Ref
        tb = add_textbox(slide, x + Inches(0.3), card_top + card_h - Inches(0.35),
                         card_w - Inches(0.6), Inches(0.3))
        set_text(tb.text_frame.paragraphs[0], ref,
                 size=9, color=BONE_3, italic=True, font=FONT_SERIF)

    # Privacy foot
    tb = add_textbox(slide, MARGIN, SLIDE_H - Inches(1.1),
                     SLIDE_W - 2 * MARGIN, Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "All processing on-device. "
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(18)
    r1.font.color.rgb = BRASS
    r1.font.italic = True
    r2 = p.add_run()
    r2.text = "No audio, no video, no GPS leaves the phone."
    r2.font.name = FONT_SERIF
    r2.font.size = Pt(18)
    r2.font.color.rgb = BONE_1

    add_chrome_bottom(slide,
                      "Module 01 · Breath ▸ Module 02 · Motion ▸ Module 03 · Heart",
                      "03 / 05")


# --- Slide 04 · Live demo -----------------------------------------------------


def fetch_qr_png(url: str) -> io.BytesIO:
    api = (
        "https://api.qrserver.com/v1/create-qr-code/"
        f"?size=900x900&ecc=M&margin=1&color=0a0b10&bgcolor=f4ece1&data={url}"
    )
    with urllib.request.urlopen(api, timeout=15) as resp:
        data = resp.read()
    return io.BytesIO(data)


def build_slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, INK_2)
    add_chrome_top(slide, "CASE 04 · LIVE DEMO", "04 / 05")

    # Left: eyebrow + headline + steps
    left_x = MARGIN
    left_w = Inches(7.0)

    tb = add_textbox(slide, left_x, Inches(1.3), left_w, Inches(0.4))
    set_text(tb.text_frame.paragraphs[0],
             "● LIVE · ON THE PROJECTOR BESIDE ME",
             size=11, color=WARN, bold=True, font=FONT_SANS)

    tb = add_textbox(slide, left_x, Inches(1.85), left_w, Inches(2.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Scan. Blow. "
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(52)
    r1.font.color.rgb = BONE_0
    r2 = p.add_run()
    r2.text = "Fill the bar."
    r2.font.name = FONT_SERIF
    r2.font.size = Pt(52)
    r2.font.color.rgb = BRASS
    r2.font.italic = True

    # Steps
    steps = [
        ("01", "Scan the code with your phone camera. You land on an onboarding form."),
        ("02", "Enter a team code."),
        ("03", "Blow into the bottom of your phone for six seconds."),
        ("04", "See your numbers, your GP letter, your team on the board."),
    ]
    step_top = Inches(4.2)
    for i, (n, t) in enumerate(steps):
        y = step_top + Inches(i * 0.55)
        # Number (serif)
        tb = add_textbox(slide, left_x, y, Inches(0.7), Inches(0.5))
        set_text(tb.text_frame.paragraphs[0], n,
                 size=22, color=BRASS, font=FONT_SERIF, italic=True)
        # Text
        tb = add_textbox(slide, left_x + Inches(0.9), y + Inches(0.08),
                         left_w - Inches(0.9), Inches(0.5))
        set_text(tb.text_frame.paragraphs[0], t,
                 size=13, color=BONE_1, font=FONT_SANS)

    # Caption box
    cap_top = Inches(6.5)
    add_rect(slide, left_x, cap_top, left_w, Inches(0.55),
             fill_color=INK_1, line_color=BRASS, line_width=Pt(1))
    tb = add_textbox(slide, left_x + Inches(0.3), cap_top + Inches(0.1),
                     left_w - Inches(0.6), Inches(0.4))
    set_text(tb.text_frame.paragraphs[0],
             "Fill the bar together. Top team takes the board.",
             size=16, color=BRASS_BRIGHT, italic=True, font=FONT_SERIF)

    # Right: QR
    qr_size = Inches(4.4)
    qr_x = SLIDE_W - MARGIN - qr_size
    qr_y = Inches(1.7)
    # Card background (white/bone)
    add_rect(slide, qr_x - Inches(0.25), qr_y - Inches(0.25),
             qr_size + Inches(0.5), qr_size + Inches(0.5),
             fill_color=BONE_0)
    # Fetch + embed QR
    demo_url = "https://mocha-cleat-chastity.ngrok-free.dev/"
    try:
        qr_stream = fetch_qr_png(demo_url)
        slide.shapes.add_picture(qr_stream, qr_x, qr_y,
                                 width=qr_size, height=qr_size)
    except Exception as err:
        # Fallback: placeholder rectangle + URL text
        add_rect(slide, qr_x, qr_y, qr_size, qr_size, fill_color=INK_0)
        tb = add_textbox(slide, qr_x, qr_y + qr_size / 2 - Inches(0.3),
                         qr_size, Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_text(p, "QR goes here\n" + str(err),
                 size=12, color=BONE_0, font=FONT_SANS, align=PP_ALIGN.CENTER)

    # URL caption
    tb = add_textbox(slide, qr_x - Inches(0.25), qr_y + qr_size + Inches(0.3),
                     qr_size + Inches(0.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p, "mocha-cleat-chastity.ngrok-free.dev",
             size=12, color=BONE_2, font=FONT_SANS, align=PP_ALIGN.CENTER)

    add_chrome_bottom(slide,
                      "GP letter written live by GLM 5.1 · no audio leaves your phone",
                      "04 / 05")


# --- Slide 05 · Business + roadmap (Q&A backup) -------------------------------


def build_slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, INK_1)
    add_chrome_top(slide, "CASE 05 · BUSINESS", "05 / 05")

    # Eyebrow
    tb = add_textbox(slide, MARGIN, Inches(1.2),
                     SLIDE_W - 2 * MARGIN, Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p, "PER-SEAT SAAS · SLACK AND TEAMS DAY ONE",
             size=12, color=BRASS, bold=True, font=FONT_SANS,
             align=PP_ALIGN.CENTER)

    # Headline
    tb = add_textbox(slide, MARGIN, Inches(1.7),
                     SLIDE_W - 2 * MARGIN, Inches(1.1))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "£3"
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(72)
    r1.font.color.rgb = BRASS
    r2 = p.add_run()
    r2.text = " a seat. Zero hardware."
    r2.font.name = FONT_SERIF
    r2.font.size = Pt(54)
    r2.font.color.rgb = BONE_0

    # Two columns: pricing + roadmap
    col_top = Inches(3.5)
    col_h = Inches(3.2)
    col_w = Inches((SLIDE_W.inches - 2 * MARGIN.inches - 0.6) / 2)
    col1_x = MARGIN
    col2_x = MARGIN + col_w + Inches(0.6)

    # -- Pricing column --
    tb = add_textbox(slide, col1_x, col_top, col_w, Inches(0.3))
    set_text(tb.text_frame.paragraphs[0], "PRICING AGAINST THE MARKET",
             size=11, color=BONE_3, bold=True, font=FONT_SANS)

    # Two price cards side by side
    card_w = Inches((col_w.inches - 0.25) / 2)
    card_h = Inches(2.0)
    card_top = col_top + Inches(0.45)

    def draw_price(x, is_us, name, cost_ccy, cost_num, per, note, note_warn=False):
        fill = INK_2 if is_us else INK_1
        border = BRASS if is_us else BONE_3
        add_rect(slide, x, card_top, card_w, card_h,
                 fill_color=fill, line_color=border,
                 line_width=Pt(1 if is_us else 0.5))
        tb = add_textbox(slide, x + Inches(0.2), card_top + Inches(0.15),
                         card_w - Inches(0.4), Inches(0.3))
        set_text(tb.text_frame.paragraphs[0], name,
                 size=10, color=BRASS if is_us else BONE_3,
                 bold=True, font=FONT_SANS)
        tb = add_textbox(slide, x + Inches(0.2), card_top + Inches(0.5),
                         card_w - Inches(0.4), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run()
        r1.text = cost_ccy
        r1.font.name = FONT_SERIF
        r1.font.size = Pt(28)
        r1.font.color.rgb = BRASS
        r2 = p.add_run()
        r2.text = cost_num
        r2.font.name = FONT_SERIF
        r2.font.size = Pt(56)
        r2.font.color.rgb = BONE_0
        r3 = p.add_run()
        r3.text = f"  {per}"
        r3.font.name = FONT_SANS
        r3.font.size = Pt(10)
        r3.font.color.rgb = BONE_3
        r3.font.bold = True
        tb = add_textbox(slide, x + Inches(0.2), card_top + card_h - Inches(0.45),
                         card_w - Inches(0.4), Inches(0.35))
        set_text(tb.text_frame.paragraphs[0], note,
                 size=10, color=WARN if note_warn else BONE_2, font=FONT_SANS)

    draw_price(col1_x, False, "FITBIT FOR WORK", "£", "12", "/SEAT/MO",
               "40% drop-off by month 3", note_warn=True)
    draw_price(col1_x + card_w + Inches(0.25), True,
               "RESONA", "£", "3", "/SEAT/MO",
               "Zero hardware · daily touchpoint")

    # Model line under pricing
    tb = add_textbox(slide, col1_x, card_top + card_h + Inches(0.25),
                     col_w, Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Quarter the cost. "
    r1.font.name = FONT_SERIF
    r1.font.size = Pt(18)
    r1.font.color.rgb = BONE_1
    r2 = p.add_run()
    r2.text = "Four times the frequency."
    r2.font.name = FONT_SERIF
    r2.font.size = Pt(18)
    r2.font.color.rgb = BRASS
    r2.font.italic = True
    p2 = tf.add_paragraph()
    r3 = p2.add_run()
    r3.text = 'Slack nudge at 2pm: "Resona check-in due."'
    r3.font.name = FONT_SERIF
    r3.font.size = Pt(14)
    r3.font.color.rgb = BONE_2
    r3.font.italic = True

    # -- Roadmap column --
    tb = add_textbox(slide, col2_x, col_top, col_w, Inches(0.3))
    set_text(tb.text_frame.paragraphs[0], "SAME PHONE · MORE BODY SYSTEMS",
             size=11, color=BONE_3, bold=True, font=FONT_SANS)

    stops = [
        ("TODAY", "Breath + Motion.", "Mic · IMU", True),
        ("Q3", "Heart.", "Camera · rPPG", False),
        ("Q4", "Sleep telemetry.", "IMU · overnight", False),
        ("2027", "Mental load.", "Voice markers", False),
    ]
    stop_top = col_top + Inches(0.45)
    for i, (when, what, modality, is_today) in enumerate(stops):
        y = stop_top + Inches(i * 0.65)
        # When
        tb = add_textbox(slide, col2_x, y, Inches(1.2), Inches(0.4))
        set_text(tb.text_frame.paragraphs[0], when,
                 size=11, color=PULSE if is_today else BRASS,
                 bold=True, font=FONT_SANS)
        # What
        tb = add_textbox(slide, col2_x + Inches(1.3), y - Inches(0.05),
                         col_w - Inches(3.5), Inches(0.5))
        set_text(tb.text_frame.paragraphs[0], what,
                 size=20, color=BRASS_BRIGHT if is_today else BONE_0,
                 font=FONT_SERIF)
        # Modality
        tb = add_textbox(slide, col2_x + col_w - Inches(2.0), y + Inches(0.05),
                         Inches(2.0), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        set_text(p, modality, size=9, color=BONE_3, bold=True,
                 font=FONT_SANS, align=PP_ALIGN.RIGHT)
        # Hairline below each stop
        if i < len(stops) - 1:
            add_rect(slide, col2_x, y + Inches(0.55),
                     col_w, Emu(9525), fill_color=BONE_3)

    # The ask bar
    ask_top = SLIDE_H - Inches(1.5)
    ask_h = Inches(0.8)
    add_rect(slide, MARGIN, ask_top, SLIDE_W - 2 * MARGIN, ask_h,
             fill_color=INK_2, line_color=BRASS, line_width=Pt(1))
    tb = add_textbox(slide, MARGIN + Inches(0.4),
                     ask_top + Inches(0.15),
                     Inches(2), Inches(0.5))
    set_text(tb.text_frame.paragraphs[0], "THE ASK",
             size=12, color=BRASS, bold=True, font=FONT_SANS)
    tb = add_textbox(slide, MARGIN + Inches(2.4),
                     ask_top + Inches(0.15),
                     SLIDE_W - 2 * MARGIN - Inches(2.8),
                     Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    set_text(p,
             "Watcha's enterprise network. Pilot with five portfolio companies, Monday.",
             size=18, color=BONE_0, font=FONT_SERIF, align=PP_ALIGN.RIGHT)

    add_chrome_bottom(slide,
                      "Every body has a rhythm. Now every team has one too.",
                      "05 / 05")


# --- Entry point --------------------------------------------------------------


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_01(prs)
    build_slide_02(prs)
    build_slide_03(prs)
    build_slide_04(prs)
    build_slide_05(prs)

    out_dir = Path(__file__).resolve().parent.parent / "submission"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "Resona.pptx"
    prs.save(str(out_path))
    print(f"wrote {out_path}  ({out_path.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
